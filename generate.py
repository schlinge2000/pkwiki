# /// script
# dependencies = ["python-pptx", "pillow", "openai", "python-dotenv", "pydantic"]
# ///
"""
Knowledge Wiki — PPTX-Generator.
Baut aus einer JSON-Outline (von synthesize.py) + POTX-Template eine fertige PPTX zusammen.

Usage:
  uv run generate.py output/demand-ai-supply-chain.json
  uv run generate.py output/demand-ai-supply-chain.json --template templates/inform-master-en.potx

Default-Template: templates/inform-master-de.potx
"""

import sys
import io
import json
import argparse
import logging
from datetime import datetime, timezone
from pathlib import Path
from typing import Optional

# UTF-8 erzwingen (Windows-Konsole)
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")
sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding="utf-8", errors="replace")

# .env laden falls vorhanden
try:
    from dotenv import load_dotenv
    load_dotenv(Path(__file__).parent / ".env")
except ImportError:
    pass

from pydantic import BaseModel
from typing import Literal

# ---------------------------------------------------------------------------
# Pfad-Konstanten
# ---------------------------------------------------------------------------

WIKI_ROOT    = Path(__file__).parent
ASSETS_DIR   = WIKI_ROOT / "assets"
OUTPUT_DIR   = WIKI_ROOT / "output"
TEMPLATES_DIR = WIKI_ROOT / "templates"
INDEX_FILE   = ASSETS_DIR / "image-index.json"

# ---------------------------------------------------------------------------
# Pydantic-Modelle (spiegeln synthesize.py wider)
# ---------------------------------------------------------------------------


class SlideOutline(BaseModel):
    nummer: int
    typ: Literal["title", "agenda", "content", "visual", "closing"]
    titel: str
    kernaussage: str
    bullet_points: list[str]
    sprecher_notiz: str
    quellen: list[str]
    bild_vorschlag: str


class PresentationOutline(BaseModel):
    titel: str
    zielgruppe: str
    sprache: Literal["de", "en"]
    folien: list[SlideOutline]
    narrative_arc: str


# ---------------------------------------------------------------------------
# Logging einrichten
# ---------------------------------------------------------------------------


def setup_logging() -> logging.Logger:
    """Richtet Console-Logging ein."""
    logger = logging.getLogger("generate")
    logger.setLevel(logging.DEBUG)
    logger.handlers.clear()

    fmt = logging.Formatter(
        fmt="%(asctime)s  %(levelname)-8s  %(message)s",
        datefmt="%Y-%m-%d %H:%M:%S",
    )
    ch = logging.StreamHandler(sys.stdout)
    ch.setLevel(logging.INFO)
    ch.setFormatter(fmt)
    logger.addHandler(ch)
    return logger


# ---------------------------------------------------------------------------
# Outline laden
# ---------------------------------------------------------------------------


def load_outline(json_path: Path) -> PresentationOutline:
    """Lädt eine PresentationOutline aus einer JSON-Datei (von synthesize.py)."""
    if not json_path.exists():
        raise FileNotFoundError(f"Outline-Datei nicht gefunden: {json_path}")
    raw = json.loads(json_path.read_text(encoding="utf-8"))
    return PresentationOutline.model_validate(raw)


# ---------------------------------------------------------------------------
# Bild-Index laden und Bild suchen
# ---------------------------------------------------------------------------


def _load_image_index(log: logging.Logger) -> list[dict]:
    """Lädt assets/image-index.json. Gibt leere Liste zurück wenn nicht vorhanden."""
    if not INDEX_FILE.exists():
        log.debug("image-index.json nicht vorhanden — keine Bild-Zuordnung möglich")
        return []
    try:
        data = json.loads(INDEX_FILE.read_text(encoding="utf-8"))
        if isinstance(data, list):
            return data
        log.warning("image-index.json hat unerwartetes Format")
        return []
    except Exception as e:
        log.warning("Fehler beim Laden von image-index.json: %s", e)
        return []


def find_best_image(bild_vorschlag: str, image_index: list[dict], log: logging.Logger) -> Optional[Path]:
    """
    Sucht im Bild-Index das semantisch passendste Bild zum Vorschlag.
    Einfaches Keyword-Matching auf dem 'beschreibung'-Feld jedes Eintrags.
    Gibt absoluten Pfad zurück oder None wenn kein passendes Bild gefunden.
    """
    if not bild_vorschlag or not image_index:
        return None

    # Suchbegriffe aus dem Vorschlag extrahieren (Kleinbuchstaben, min. 3 Zeichen)
    keywords = [w.lower() for w in bild_vorschlag.split() if len(w) >= 3]
    if not keywords:
        return None

    bestes_ergebnis: Optional[dict] = None
    bester_score = 0

    for eintrag in image_index:
        beschreibung = (eintrag.get("beschreibung") or "").lower()
        # Score = Anzahl Treffer der Keywords in der Beschreibung
        score = sum(1 for kw in keywords if kw in beschreibung)
        if score > bester_score:
            bester_score = score
            bestes_ergebnis = eintrag

    if bestes_ergebnis is None or bester_score == 0:
        log.debug("Kein passendes Bild für '%s' gefunden", bild_vorschlag)
        return None

    # Pfad auflösen — image_path ist relativ zu WIKI_ROOT
    image_path_str = bestes_ergebnis.get("image_path", "")
    if not image_path_str:
        return None

    abs_path = WIKI_ROOT / image_path_str
    if abs_path.exists():
        log.debug(
            "Bild für '%s' gefunden: %s (Score: %d)",
            bild_vorschlag,
            image_path_str,
            bester_score,
        )
        return abs_path

    log.debug("Bild-Pfad existiert nicht: %s", abs_path)
    return None


# ---------------------------------------------------------------------------
# Slide einfügen
# ---------------------------------------------------------------------------


def _get_slide_layout(prs, typ: str, log: logging.Logger):
    """
    Wählt das passende Slide-Layout basierend auf dem Folientyp.
    Fällt auf verfügbare Layouts zurück wenn das gewünschte nicht existiert.
    """
    # Mapping: Folientyp → bevorzugter Layout-Index
    layout_map = {
        "title":   0,
        "agenda":  1,
        "content": 2,
        "visual":  6,
        "closing": 0,
    }
    preferred_idx = layout_map.get(typ, 2)
    max_idx = len(prs.slide_layouts) - 1

    if preferred_idx <= max_idx:
        return prs.slide_layouts[preferred_idx]

    # Fallback: letztes verfügbares Layout
    log.warning(
        "Layout-Index %d nicht verfügbar (Template hat %d Layouts) — nutze Index %d",
        preferred_idx,
        max_idx + 1,
        max_idx,
    )
    return prs.slide_layouts[max_idx]


def _befuelle_platzhalter(slide, slide_outline: SlideOutline, log: logging.Logger) -> None:
    """
    Befüllt Titel-, Bullet- und Notiz-Platzhalter einer Folie.
    Nicht vorhandene Platzhalter werden übersprungen (kein Absturz).
    """
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN

    # Platzhalter nach Typ gruppieren
    titel_ph = None
    bullet_ph = None

    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:
            # Index 0 = Titel-Platzhalter
            titel_ph = ph
        elif idx == 1 and bullet_ph is None:
            # Index 1 = Erster Inhalts-Platzhalter (Bullets)
            bullet_ph = ph

    # Titel befüllen
    if titel_ph is not None:
        try:
            titel_ph.text = slide_outline.titel
        except Exception as e:
            log.warning("Folie %d: Titel konnte nicht gesetzt werden: %s", slide_outline.nummer, e)
    else:
        log.debug("Folie %d: Kein Titel-Platzhalter (idx=0) im Layout", slide_outline.nummer)

    # Bullet-Points befüllen
    if bullet_ph is not None and slide_outline.bullet_points:
        try:
            tf = bullet_ph.text_frame
            tf.clear()
            for i, bullet in enumerate(slide_outline.bullet_points):
                if i == 0:
                    # Ersten Paragraph direkt nutzen
                    para = tf.paragraphs[0]
                else:
                    para = tf.add_paragraph()
                para.text = bullet
        except Exception as e:
            log.warning(
                "Folie %d: Bullet-Points konnten nicht gesetzt werden: %s",
                slide_outline.nummer,
                e,
            )
    elif bullet_ph is None and slide_outline.bullet_points:
        log.debug("Folie %d: Kein Inhalts-Platzhalter (idx=1) im Layout", slide_outline.nummer)

    # Speaker Notes befüllen
    if slide_outline.sprecher_notiz:
        try:
            notes_slide = slide.notes_slide
            notes_tf = notes_slide.notes_text_frame
            notes_tf.text = slide_outline.sprecher_notiz
        except Exception as e:
            log.warning("Folie %d: Speaker Notes konnten nicht gesetzt werden: %s", slide_outline.nummer, e)


def _fuege_bild_ein(slide, image_path: Path, log: logging.Logger) -> None:
    """
    Fügt ein Bild in einen vorhandenen Bild-Platzhalter ein.
    Wenn kein Bild-Platzhalter vorhanden, wird das Bild übersprungen.
    """
    from pptx.enum.shapes import PP_PLACEHOLDER

    # Bild-Platzhalter suchen (idx >= 2 oder typ OBJECT/PICTURE)
    bild_ph = None
    for ph in slide.placeholders:
        ph_type = ph.placeholder_format.type
        # Bild- oder Objekt-Platzhalter
        if ph_type in (
            PP_PLACEHOLDER.OBJECT,
            PP_PLACEHOLDER.PICTURE,
            PP_PLACEHOLDER.BODY,
        ):
            # Nur Platzhalter mit idx >= 2 (kein Titel/Content)
            if ph.placeholder_format.idx >= 2:
                bild_ph = ph
                break

    if bild_ph is None:
        log.debug(
            "Kein Bild-Platzhalter im Layout — Bild '%s' wird nicht eingefügt", image_path.name
        )
        return

    try:
        bild_ph.insert_picture(str(image_path))
        log.debug("Bild eingefügt: %s", image_path.name)
    except Exception as e:
        log.warning("Bild konnte nicht eingefügt werden (%s): %s", image_path.name, e)


def add_slide(prs, slide_outline: SlideOutline, image_path: Optional[Path], log: logging.Logger):
    """
    Fügt eine neue Folie zur Präsentation hinzu.
    Layout wird anhand des Folientyps gewählt.
    Platzhalter werden befüllt, Bild optional eingefügt.
    Gibt die hinzugefügte Folie zurück.
    """
    layout = _get_slide_layout(prs, slide_outline.typ, log)

    # Folie aus Layout erzeugen
    slide = prs.slides.add_slide(layout)

    # Platzhalter befüllen
    _befuelle_platzhalter(slide, slide_outline, log)

    # Bild einfügen wenn vorhanden
    if image_path is not None:
        _fuege_bild_ein(slide, image_path, log)

    return slide


# ---------------------------------------------------------------------------
# Hauptfunktion
# ---------------------------------------------------------------------------


def generate(outline_path: Path, template_path: Path, log: logging.Logger) -> Path:
    """
    Hauptfunktion: Template öffnen, alle vorhandenen Slides entfernen,
    für jede Folie aus der Outline add_slide() aufrufen, PPTX speichern.
    Gibt den Pfad der gespeicherten PPTX zurück.
    """
    from pptx import Presentation

    # Template prüfen
    if not template_path.exists():
        raise FileNotFoundError(
            f"Template nicht gefunden: {template_path}\n"
            f"Bitte POTX-Template in '{TEMPLATES_DIR}/' ablegen.\n"
            f"Verfügbare Templates: {[t.name for t in TEMPLATES_DIR.glob('*.potx')]}"
        )

    log.info("Template geladen: %s", template_path.name)
    prs = Presentation(str(template_path))

    # Alle vorhandenen Folien entfernen — leere Basis
    # pptx speichert Slides in einer XML-Liste; wir löschen alle rxml-Elemente
    xml_slides = prs.slides._sldIdLst
    for sld_id in list(xml_slides):
        xml_slides.remove(sld_id)
    log.info("Vorhandene Template-Folien entfernt (leere Basis)")

    # Outline laden
    outline = load_outline(outline_path)
    log.info(
        "Outline geladen: '%s' — %d Folien | Zielgruppe: %s | Sprache: %s",
        outline.titel,
        len(outline.folien),
        outline.zielgruppe,
        outline.sprache,
    )

    # Bild-Index laden
    image_index = _load_image_index(log)
    log.info("Bild-Index geladen: %d Einträge", len(image_index))

    # Folien aufbauen
    for folie in outline.folien:
        # Passendes Bild suchen
        image_path = find_best_image(folie.bild_vorschlag, image_index, log)

        add_slide(prs, folie, image_path, log)
        log.info(
            "  Folie %2d/%d  [%-8s]  %s%s",
            folie.nummer,
            len(outline.folien),
            folie.typ,
            folie.titel[:50],
            f"  → Bild: {image_path.name}" if image_path else "",
        )

    # Ausgabepfad bestimmen: output/<slug>-<timestamp>.pptx
    # Slug aus dem JSON-Dateinamen ableiten
    slug = outline_path.stem
    timestamp = datetime.now(timezone.utc).strftime("%Y%m%d-%H%M%S")
    OUTPUT_DIR.mkdir(exist_ok=True)
    output_path = OUTPUT_DIR / f"{slug}-{timestamp}.pptx"

    prs.save(str(output_path))
    log.info("PPTX gespeichert: %s", output_path)

    return output_path


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------


def parse_args() -> argparse.Namespace:
    """Parst CLI-Argumente."""
    parser = argparse.ArgumentParser(
        description="Baut aus einer JSON-Outline + POTX-Template eine fertige PPTX zusammen.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""Beispiele:
  uv run generate.py output/demand-ai-supply-chain.json
  uv run generate.py output/demand-ai-supply-chain.json --template templates/inform-master-en.potx
""",
    )
    parser.add_argument(
        "outline",
        type=Path,
        help="Pfad zur JSON-Outline (Ausgabe von synthesize.py)",
    )
    parser.add_argument(
        "--template",
        type=Path,
        default=TEMPLATES_DIR / "inform-master-de.potx",
        help="Pfad zum POTX-Template (default: templates/inform-master-de.potx)",
    )
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    log = setup_logging()

    # Outline-Pfad auflösen (relativ zum CWD oder absolut)
    outline_path = args.outline.resolve()

    try:
        output_path = generate(outline_path, args.template, log)
        print(f"\nErfolgreich erstellt: {output_path}")
    except FileNotFoundError as e:
        print(f"\nFEHLER: {e}", file=sys.stderr)
        sys.exit(1)
    except Exception as e:
        log.exception("Unerwarteter Fehler: %s", e)
        sys.exit(1)


if __name__ == "__main__":
    main()
