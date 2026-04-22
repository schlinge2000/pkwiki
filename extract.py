# /// script
# dependencies = ["openai", "python-pptx", "python-docx", "pdfplumber", "pymupdf", "pillow", "python-dotenv"]
# ///
"""
Knowledge Wiki — Dokumentenextraktion mit Azure OpenAI Vision.
Usage: uv run extract.py <dateipfad>   (PPTX, DOCX oder PDF)

Benoetigt in .env:
  AZURE_OPENAI_API_KEY=...
  AZURE_OPENAI_ENDPOINT=https://<resource>.openai.azure.com/
  AZURE_OPENAI_DEPLOYMENT=gpt-4o
"""

import sys
import io
import os
import time
import base64
from pathlib import Path
from concurrent.futures import ThreadPoolExecutor, as_completed

# UTF-8 erzwingen (Windows-Konsole)
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")
sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding="utf-8", errors="replace")

# .env laden falls vorhanden
try:
    from dotenv import load_dotenv
    load_dotenv(Path(__file__).parent / ".env")
except ImportError:
    pass

WIKI_ROOT = Path(__file__).parent
RAW_DIR = WIKI_ROOT / "raw"
CACHE_DIR = RAW_DIR / ".cache"

# Bildverarbeitungs-Konstanten
MAX_IMAGE_WIDTH = 800   # Maximale Breite fuer Vision-API (Pixel)
JPEG_QUALITY    = 75    # JPEG-Kompressionsqualitaet
MAX_WORKERS     = 3     # Parallele Folien-Verarbeitungs-Threads


def to_jpeg_bytes(image) -> bytes | None:
    """Konvertiert ein PIL-Image zu komprimiertem JPEG (max. MAX_IMAGE_WIDTH Breite)."""
    try:
        from PIL import Image
        if not isinstance(image, Image.Image):
            return None
        # Groesse begrenzen um API-Kosten zu senken
        if image.width > MAX_IMAGE_WIDTH:
            ratio = MAX_IMAGE_WIDTH / image.width
            new_size = (MAX_IMAGE_WIDTH, int(image.height * ratio))
            image = image.resize(new_size, Image.LANCZOS)
        # JPEG unterstuetzt kein RGBA — auf RGB konvertieren
        if image.mode != "RGB":
            image = image.convert("RGB")
        buf = io.BytesIO()
        image.save(buf, format="JPEG", quality=JPEG_QUALITY)
        return buf.getvalue()
    except Exception as e:
        print(f"  Bildkonvertierung fehlgeschlagen: {e}", file=sys.stderr)
        return None


def _vision_api_call_once(images_b64: list[str], slide_context: str, slide_num: int) -> list[str]:
    """Einzelner Vision-API-Aufruf fuer alle Bilder einer Folie."""
    from openai import AzureOpenAI
    client = AzureOpenAI(
        api_key=os.environ["AZURE_OPENAI_API_KEY"],
        azure_endpoint=os.environ["AZURE_OPENAI_ENDPOINT"],
        api_version=os.environ.get("AZURE_OPENAI_API_VERSION", "2025-04-01-preview"),
    )
    deployment = os.environ.get("AZURE_OPENAI_DEPLOYMENT", "gpt-4o")

    # Prompt: fachlicher Fokus, Deutsch, eine Zeile pro Bild
    content: list[dict] = [
        {
            "type": "text",
            "text": (
                f"Folie {slide_num} einer Praesentation.\n"
                f"Textinhalt: {slide_context[:500]}\n\n"
                "Beschreibe praegnant auf Deutsch was auf den Bildern zu sehen ist. "
                "Fokus auf fachlichen Inhalt (Diagramme, Kennzahlen, Prozesse, Architektur). "
                "Je Bild eine Zeile, keine Einleitung."
            )
        }
    ]
    for b64 in images_b64:
        content.append({
            "type": "image_url",
            "image_url": {"url": f"data:image/jpeg;base64,{b64}", "detail": "low"}
        })

    response = client.chat.completions.create(
        model=deployment,
        messages=[{"role": "user", "content": content}],
        max_completion_tokens=500,
    )
    result_text = (response.choices[0].message.content or "").strip()
    return [line.strip() for line in result_text.splitlines() if line.strip()]


class VisionAPIError(Exception):
    """Nicht-behebarer Vision-API-Fehler (Konnektivität, Auth, etc.)."""


def describe_slide_images(images_b64: list[str], slide_context: str, slide_num: int) -> list[str]:
    """Beschreibt Bilder einer Folie via Vision-API mit Retry bei Rate-Limit (max. 4 Versuche).

    Raises VisionAPIError bei Konnektivitäts- oder Auth-Fehlern (nicht bei Rate-Limit).
    """
    if not images_b64:
        return []
    for attempt in range(4):
        try:
            return _vision_api_call_once(images_b64, slide_context, slide_num)
        except Exception as e:
            err = str(e)
            if "429" in err or "too_many_requests" in err.lower() or "rate" in err.lower():
                wait = 10 * (attempt + 1)
                print(f"  Rate-Limit (Folie {slide_num}), warte {wait}s...", flush=True)
                time.sleep(wait)
            else:
                # Konnektivitätsfehler, Auth-Fehler o.ä. — nicht still schlucken
                raise VisionAPIError(f"Folie {slide_num}: {e}") from e
    raise VisionAPIError(f"Folie {slide_num}: Rate-Limit nach 4 Versuchen")


def _process_slide(args: tuple) -> tuple[int, list[str], list[str]]:
    """Verarbeitet eine Folie: Text extrahieren + Bilder beschreiben (ThreadPoolExecutor-kompatibel).

    Raises VisionAPIError bei Konnektivitätsproblemen (propagiert aus describe_slide_images).
    """
    slide_idx, slide, use_vision = args
    slide_num = slide_idx + 1
    text_lines: list[str] = []
    images_b64: list[str] = []

    # Textinhalte aller Shapes sammeln
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    text_lines.append(text)

    # Bilder zu JPEG konvertieren (nur wenn Vision aktiv)
    if use_vision:
        for shape in slide.shapes:
            try:
                # hasattr allein reicht nicht — shape.image kann ValueError werfen
                # wenn das Bild verknüpft (nicht eingebettet) ist
                if not hasattr(shape, "image"):
                    continue
                _ = shape.image  # Probe-Zugriff um ValueError frühzeitig zu fangen
                from PIL import Image
                img = Image.open(io.BytesIO(shape.image.blob))
                jpeg = to_jpeg_bytes(img)
                if jpeg:
                    images_b64.append(base64.b64encode(jpeg).decode("ascii"))
            except Exception:
                # EMF/WMF-Vektorgrafiken koennen von PIL nicht geoefffnet werden — ueberspringen
                pass

    # Vision-API aufrufen falls Bilder vorhanden — VisionAPIError propagiert nach oben
    slide_context = " ".join(text_lines)
    image_descs = describe_slide_images(images_b64, slide_context, slide_num) if images_b64 else []

    return slide_num, text_lines, image_descs


MAX_VISION_ERRORS = 3  # Ab dieser Anzahl Verbindungsfehler wird abgebrochen


def extract_pptx(path: Path) -> str:
    """Extrahiert PPTX: Text + Vision-Bildbeschreibungen, parallel pro Folie.

    Raises RuntimeError wenn zu viele Vision-API-Fehler auftreten (Konnektivitätsverlust).
    """
    from pptx import Presentation

    use_vision = bool(os.environ.get("AZURE_OPENAI_API_KEY"))
    prs = Presentation(path)
    total = len(prs.slides)
    print(f"  {total} Folien | Vision: {'aktiv' if use_vision else 'inaktiv'}", flush=True)

    results: dict[int, tuple[list[str], list[str]]] = {}
    vision_errors = 0

    with ThreadPoolExecutor(max_workers=MAX_WORKERS) as executor:
        args_list = [(i, slide, use_vision) for i, slide in enumerate(prs.slides)]
        futures = {executor.submit(_process_slide, args): args[0] for args in args_list}
        for future in as_completed(futures):
            try:
                slide_num, text_lines, image_descs = future.result()
                results[slide_num] = (text_lines, image_descs)
                print(f"  Folie {slide_num}/{total} ({len(image_descs)} Bildbeschreibungen)", flush=True)
            except VisionAPIError as e:
                vision_errors += 1
                slide_idx = futures[future]
                print(f"  FEHLER Folie {slide_idx + 1}: {e}", file=sys.stderr)
                results[slide_idx + 1] = ([], [])
                if vision_errors >= MAX_VISION_ERRORS:
                    raise RuntimeError(
                        f"Extraktion abgebrochen — {vision_errors} Vision-API-Fehler in Folge. "
                        "Konnektivitätsproblem? Cache wird nicht geschrieben."
                    ) from e

    # Ausgabe in Folien-Reihenfolge zusammensetzen
    lines = [f"# {path.stem}\n", f"**Quelle:** {path.name}\n"]
    for slide_num in sorted(results):
        text_lines, image_descs = results[slide_num]
        if not text_lines and not image_descs:
            continue
        lines.append(f"\n## Folie {slide_num}\n")
        lines.extend(text_lines)
        if image_descs:
            lines.append("\n**Bilder:**")
            for desc in image_descs:
                lines.append(f"- {desc}")

    return "\n".join(lines)


def extract_docx(path: Path) -> str:
    """Extrahiert Text aus DOCX unter Beibehaltung der Ueberschriftenstruktur."""
    from docx import Document

    doc = Document(path)
    lines = [f"# {path.stem}\n", f"**Quelle:** {path.name}\n"]

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style = para.style.name if para.style else ""
        if "Heading 1" in style:
            lines.append(f"\n# {text}")
        elif "Heading 2" in style:
            lines.append(f"\n## {text}")
        elif "Heading 3" in style:
            lines.append(f"\n### {text}")
        else:
            lines.append(text)

    return "\n".join(lines)


def _extract_pdf_pdfplumber(path: Path) -> str:
    """PDF-Extraktion mit pdfplumber (Standardweg)."""
    import pdfplumber

    lines = [f"# {path.stem}\n", f"**Quelle:** {path.name}\n"]
    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            text = (page.extract_text() or "").strip()
            if text:
                lines.append(f"\n## Seite {i}\n")
                lines.append(text)
    return "\n".join(lines)


def _extract_pdf_pymupdf(path: Path) -> str:
    """PDF-Extraktion mit PyMuPDF — robuster Fallback fuer beschaedigte PDFs."""
    import fitz  # PyMuPDF

    lines = [f"# {path.stem}\n", f"**Quelle:** {path.name}\n"]
    doc = fitz.open(str(path))
    for i, page in enumerate(doc, 1):
        text = page.get_text().strip()
        if text:
            lines.append(f"\n## Seite {i}\n")
            lines.append(text)
    doc.close()
    return "\n".join(lines)


def extract_pdf(path: Path) -> str:
    """PDF-Extraktion: erst pdfplumber, bei Fehler automatisch Fallback auf PyMuPDF."""
    try:
        result = _extract_pdf_pdfplumber(path)
        print("  pdfplumber erfolgreich", flush=True)
        return result
    except Exception as e:
        print(f"  pdfplumber fehlgeschlagen ({e.__class__.__name__}) — versuche PyMuPDF...", flush=True)
        result = _extract_pdf_pymupdf(path)
        print("  PyMuPDF erfolgreich (Fallback)", flush=True)
        return result


def extract(file_path: str) -> Path:
    """Hauptfunktion: Dokument extrahieren und als .md im Cache ablegen."""
    path = Path(file_path).resolve()
    suffix = path.suffix.lower()

    print(f"Extrahiere: {path.name}", flush=True)

    try:
        if suffix == ".pptx":
            content = extract_pptx(path)
        elif suffix == ".docx":
            content = extract_docx(path)
        elif suffix == ".pdf":
            content = extract_pdf(path)
        else:
            print(f"Kein Extraktor fuer Dateityp '{suffix}' — ueberspringe.")
            sys.exit(0)
    except RuntimeError as e:
        print(f"FEHLER: {e}", file=sys.stderr)
        sys.exit(1)

    # Cache-Pfad: raw/.cache/<unterordner>/<dateiname>.md
    raw_root = path.parent
    while raw_root.name != "raw" and raw_root.parent != raw_root:
        raw_root = raw_root.parent
    rel = path.relative_to(raw_root)
    out_path = raw_root / ".cache" / rel.with_suffix(".md")
    out_path.parent.mkdir(parents=True, exist_ok=True)
    out_path.write_text(content, encoding="utf-8")
    print(f"Gespeichert: {out_path.relative_to(raw_root.parent)}", flush=True)
    return out_path


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: uv run extract.py <dateipfad>")
        sys.exit(1)
    extract(sys.argv[1])
