# /// script
# dependencies = ["openai", "python-pptx", "python-docx", "pdfplumber", "pymupdf", "pillow", "python-dotenv"]
# ///
"""
Knowledge Wiki — Granulare Bildextraktion aus PPTX/DOCX/PDF mit Vision-API-Beschreibungen.

Usage:
  uv run extract-images.py                      # alle Dateien in raw/
  uv run extract-images.py raw/slides/foo.pptx  # einzelne Datei
  uv run extract-images.py --index-only         # nur Image-Index neu aufbauen

Benoetigt in .env:
  AZURE_OPENAI_API_KEY=...
  AZURE_OPENAI_ENDPOINT=https://<resource>.openai.azure.com/
  AZURE_OPENAI_DEPLOYMENT=gpt-4o
"""

import sys
import io
import os
import json
import time
import base64
from pathlib import Path
from datetime import datetime, timezone

# UTF-8 erzwingen (Windows-Konsole)
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")
sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding="utf-8", errors="replace")

# .env laden falls vorhanden
try:
    from dotenv import load_dotenv
    load_dotenv(Path(__file__).parent / ".env")
except ImportError:
    pass

# Verzeichnis-Konstanten
WIKI_ROOT   = Path(__file__).parent
RAW_DIR     = WIKI_ROOT / "raw"
CACHE_DIR   = RAW_DIR / ".cache" / ".images"
ASSETS_DIR  = WIKI_ROOT / "assets"
INDEX_FILE  = ASSETS_DIR / "image-index.json"

# Bildverarbeitungs-Konstanten
MAX_IMAGE_WIDTH = 1200   # Maximale Breite beim Speichern (Pixel)
JPEG_QUALITY    = 85     # JPEG-Kompressionsqualitaet
# Kleinere Groesse fuer den Vision-API-Call (spart Kosten)
VISION_MAX_WIDTH = 800


# ---------------------------------------------------------------------------
# Hilfsfunktionen: Bildkonvertierung
# ---------------------------------------------------------------------------

def _to_jpeg_bytes(image, max_width: int = MAX_IMAGE_WIDTH, quality: int = JPEG_QUALITY) -> bytes | None:
    """Konvertiert ein PIL-Image zu JPEG (max. max_width Breite). Gibt None bei Fehler zurueck."""
    try:
        from PIL import Image
        if not isinstance(image, Image.Image):
            return None
        # Groesse begrenzen
        if image.width > max_width:
            ratio = max_width / image.width
            new_size = (max_width, int(image.height * ratio))
            image = image.resize(new_size, Image.LANCZOS)
        # JPEG benoetigt RGB-Modus
        if image.mode != "RGB":
            image = image.convert("RGB")
        buf = io.BytesIO()
        image.save(buf, format="JPEG", quality=quality)
        return buf.getvalue()
    except Exception as e:
        print(f"  Bildkonvertierung fehlgeschlagen: {e}", file=sys.stderr)
        return None


def _blob_to_pil(blob: bytes):
    """Oeffnet einen Rohbyte-Blob als PIL-Image. Gibt None zurueck wenn nicht lesbar (z.B. EMF/WMF)."""
    try:
        from PIL import Image
        return Image.open(io.BytesIO(blob))
    except Exception:
        return None


def _pil_size(image) -> tuple[int, int]:
    """Gibt (width, height) eines PIL-Images zurueck."""
    return image.width, image.height


# ---------------------------------------------------------------------------
# Vision-API
# ---------------------------------------------------------------------------

def describe_image(jpeg_bytes: bytes, context: str = "") -> str:
    """
    Ruft die Azure OpenAI Vision-API auf und gibt eine deutsche Bildbeschreibung zurueck.
    Bei Fehler: Fehlertext als Beschreibung (kein Absturz).
    """
    if not os.environ.get("AZURE_OPENAI_API_KEY"):
        return "[Vision-API nicht konfiguriert]"

    # Kleineres Bild fuer den API-Aufruf erstellen
    try:
        from PIL import Image
        img = Image.open(io.BytesIO(jpeg_bytes))
        vision_bytes = _to_jpeg_bytes(img, max_width=VISION_MAX_WIDTH, quality=75)
        if not vision_bytes:
            vision_bytes = jpeg_bytes
    except Exception:
        vision_bytes = jpeg_bytes

    b64 = base64.b64encode(vision_bytes).decode("ascii")

    for attempt in range(4):
        try:
            from openai import AzureOpenAI
            client = AzureOpenAI(
                api_key=os.environ["AZURE_OPENAI_API_KEY"],
                azure_endpoint=os.environ["AZURE_OPENAI_ENDPOINT"],
                api_version=os.environ.get("AZURE_OPENAI_API_VERSION", "2025-04-01-preview"),
            )
            deployment = os.environ.get("AZURE_OPENAI_DEPLOYMENT", "gpt-4o")

            content: list[dict] = [
                {
                    "type": "text",
                    "text": (
                        f"{'Kontext: ' + context[:300] if context else ''}\n\n"
                        "Beschreibe praegnant auf Deutsch was auf dem Bild zu sehen ist. "
                        "Fokus auf fachlichen Inhalt (Diagramme, Kennzahlen, Prozesse, Architektur, Grafiken). "
                        "Maximal 3 Saetze, keine Einleitung."
                    ).strip(),
                },
                {
                    "type": "image_url",
                    "image_url": {"url": f"data:image/jpeg;base64,{b64}", "detail": "low"},
                },
            ]

            response = client.chat.completions.create(
                model=deployment,
                messages=[{"role": "user", "content": content}],
                max_completion_tokens=300,
            )
            return (response.choices[0].message.content or "").strip()

        except Exception as e:
            err = str(e)
            if "429" in err or "too_many_requests" in err.lower() or "rate" in err.lower():
                wait = 10 * (attempt + 1)
                print(f"    Rate-Limit, warte {wait}s...", flush=True)
                time.sleep(wait)
            else:
                print(f"    Vision-Fehler: {e}", file=sys.stderr)
                return f"[Vision-Fehler: {e}]"

    return "[Vision-Fehler: Rate-Limit nach 4 Versuchen]"


# ---------------------------------------------------------------------------
# Meta-JSON: Schreiben und Lesen
# ---------------------------------------------------------------------------

def _write_meta(meta_path: Path, meta: dict) -> None:
    """Schreibt eine .meta.json-Datei."""
    meta_path.write_text(json.dumps(meta, ensure_ascii=False, indent=2), encoding="utf-8")


def _read_meta(meta_path: Path) -> dict | None:
    """Liest eine .meta.json-Datei. Gibt None zurueck wenn nicht lesbar."""
    try:
        return json.loads(meta_path.read_text(encoding="utf-8"))
    except Exception:
        return None


def _already_processed(meta_path: Path) -> bool:
    """Prueft ob eine .meta.json bereits existiert (Skip-Logik)."""
    return meta_path.exists()


# ---------------------------------------------------------------------------
# Pfad-Hilfsfunktionen
# ---------------------------------------------------------------------------

def _source_type(path: Path) -> str:
    """Gibt den Quelltyp ('pptx', 'docx', 'pdf') zurueck."""
    return path.suffix.lower().lstrip(".")


def _image_dir(source_path: Path) -> Path:
    """
    Gibt das Ausgabeverzeichnis fuer Bilder einer Quelldatei zurueck.
    raw/slides/foo.pptx  →  raw/.cache/.images/slides/foo/
    raw/pdfs/bar.pdf     →  raw/.cache/.images/pdfs/bar/
    """
    rel = source_path.relative_to(RAW_DIR)   # z.B. slides/foo.pptx
    subdir = rel.parent                        # z.B. slides
    stem = source_path.stem                    # z.B. foo
    return CACHE_DIR / subdir / stem


def _rel_to_wiki(path: Path) -> str:
    """Gibt den Pfad relativ zum WIKI_ROOT als Forward-Slash-String zurueck."""
    try:
        return path.relative_to(WIKI_ROOT).as_posix()
    except ValueError:
        return str(path)


# ---------------------------------------------------------------------------
# PPTX-Extraktion
# ---------------------------------------------------------------------------

def extract_images_pptx(path: Path) -> list[dict]:
    """
    Iteriert alle Slides und alle Shapes mit .image.
    Speichert jedes Bild als JPEG + .meta.json, ruft Vision-API auf.
    Gibt Liste aller Meta-Dicts zurueck.
    """
    from pptx import Presentation

    out_dir = _image_dir(path)
    out_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation(path)
    total_slides = len(prs.slides)
    print(f"  PPTX: {total_slides} Folien", flush=True)

    metas: list[dict] = []

    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1

        # Folientext als Kontext fuer Vision-API sammeln
        context_parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        context_parts.append(t)
        slide_context = " ".join(context_parts)

        shape_img_idx = 0
        for shape in slide.shapes:
            # Blob vorab lesen — hasattr() genuegt nicht, da python-pptx bei
            # verlinkten Bildern ValueError statt AttributeError wirft.
            try:
                blob = shape.image.blob  # wirft ValueError wenn kein eingebettetes Bild
                if not blob:
                    continue
            except (AttributeError, ValueError, KeyError):
                # Kein Bild, verlinktes Bild oder GroupShape — ueberspringen
                continue
            except Exception:
                continue

            shape_img_idx += 1
            # Dateiname: f<folie>-s<shape>.jpg
            img_filename = f"f{slide_num}-s{shape_img_idx}.jpg"
            img_path = out_dir / img_filename
            meta_path = out_dir / f"{img_filename}.meta.json"

            # Skip wenn bereits verarbeitet
            if _already_processed(meta_path):
                print(f"    Folie {slide_num} Shape {shape_img_idx}: bereits verarbeitet — ueberspringe", flush=True)
                existing = _read_meta(meta_path)
                if existing:
                    metas.append(existing)
                continue

            # Bild aus Shape lesen
            try:
                from PIL import Image
                # Palette-Bilder mit Transparency explizit nach RGBA konvertieren (PIL-Warning vermeiden)
                pil_img = Image.open(io.BytesIO(blob))
                if pil_img.mode == "P" and "transparency" in pil_img.info:
                    pil_img = pil_img.convert("RGBA")
                jpeg_bytes = _to_jpeg_bytes(pil_img)
                if not jpeg_bytes:
                    print(f"    Folie {slide_num} Shape {shape_img_idx}: Konvertierung fehlgeschlagen — ueberspringe", file=sys.stderr)
                    continue
                width_px, height_px = _pil_size(Image.open(io.BytesIO(jpeg_bytes)))
            except Exception as e:
                # EMF/WMF-Vektorgrafiken koennen nicht geoeffnet werden
                print(f"    Folie {slide_num} Shape {shape_img_idx}: nicht lesbar ({e.__class__.__name__}: {e}) — ueberspringe", file=sys.stderr)
                continue

            # JPEG speichern
            img_path.write_bytes(jpeg_bytes)

            # Vision-API aufrufen
            print(f"    Folie {slide_num} Shape {shape_img_idx}: Vision-API...", flush=True)
            beschreibung = describe_image(jpeg_bytes, context=slide_context)

            # Meta speichern
            meta = {
                "source_file": _rel_to_wiki(path),
                "source_type": "pptx",
                "folie": slide_num,
                "shape_index": shape_img_idx,
                "image_path": _rel_to_wiki(img_path),
                "beschreibung": beschreibung,
                "width_px": width_px,
                "height_px": height_px,
                "extracted_at": datetime.now(timezone.utc).isoformat(),
            }
            _write_meta(meta_path, meta)
            metas.append(meta)
            print(f"    Folie {slide_num} Shape {shape_img_idx}: gespeichert ({width_px}x{height_px})", flush=True)

    return metas


# ---------------------------------------------------------------------------
# PDF-Extraktion
# ---------------------------------------------------------------------------

def extract_images_pdf(path: Path) -> list[dict]:
    """
    Nutzt PyMuPDF page.get_images() um alle eingebetteten Bilder zu extrahieren.
    Speichert jedes Bild als JPEG + .meta.json, ruft Vision-API auf.
    """
    import fitz  # PyMuPDF

    out_dir = _image_dir(path)
    out_dir.mkdir(parents=True, exist_ok=True)

    doc = fitz.open(str(path))
    total_pages = doc.page_count
    print(f"  PDF: {total_pages} Seiten", flush=True)

    metas: list[dict] = []

    for page_num_0, page in enumerate(doc):
        page_num = page_num_0 + 1
        images = page.get_images(full=True)

        for img_idx, img_info in enumerate(images, 1):
            xref = img_info[0]
            img_filename = f"p{page_num}-img{img_idx}.jpg"
            img_path = out_dir / img_filename
            meta_path = out_dir / f"{img_filename}.meta.json"

            # Skip wenn bereits verarbeitet
            if _already_processed(meta_path):
                print(f"    Seite {page_num} Bild {img_idx}: bereits verarbeitet — ueberspringe", flush=True)
                existing = _read_meta(meta_path)
                if existing:
                    metas.append(existing)
                continue

            # Bilddaten aus PDF extrahieren
            try:
                base_img = doc.extract_image(xref)
                img_bytes = base_img["image"]
                from PIL import Image
                pil_img = Image.open(io.BytesIO(img_bytes))
                jpeg_bytes = _to_jpeg_bytes(pil_img)
                if not jpeg_bytes:
                    print(f"    Seite {page_num} Bild {img_idx}: Konvertierung fehlgeschlagen — ueberspringe", file=sys.stderr)
                    continue
                width_px, height_px = _pil_size(Image.open(io.BytesIO(jpeg_bytes)))
            except Exception as e:
                print(f"    Seite {page_num} Bild {img_idx}: nicht lesbar ({e.__class__.__name__}) — ueberspringe", file=sys.stderr)
                continue

            # Sehr kleine Bilder (Icons, Dekoelemente) ueberspringen
            if width_px < 50 or height_px < 50:
                print(f"    Seite {page_num} Bild {img_idx}: zu klein ({width_px}x{height_px}) — ueberspringe", flush=True)
                continue

            # JPEG speichern
            img_path.write_bytes(jpeg_bytes)

            # Seitentext als Kontext
            page_context = page.get_text().strip()[:400]

            # Vision-API aufrufen
            print(f"    Seite {page_num} Bild {img_idx}: Vision-API...", flush=True)
            beschreibung = describe_image(jpeg_bytes, context=page_context)

            # Meta speichern
            meta = {
                "source_file": _rel_to_wiki(path),
                "source_type": "pdf",
                "seite": page_num,
                "bild_index": img_idx,
                "image_path": _rel_to_wiki(img_path),
                "beschreibung": beschreibung,
                "width_px": width_px,
                "height_px": height_px,
                "extracted_at": datetime.now(timezone.utc).isoformat(),
            }
            _write_meta(meta_path, meta)
            metas.append(meta)
            print(f"    Seite {page_num} Bild {img_idx}: gespeichert ({width_px}x{height_px})", flush=True)

    doc.close()
    return metas


# ---------------------------------------------------------------------------
# DOCX-Extraktion
# ---------------------------------------------------------------------------

def extract_images_docx(path: Path) -> list[dict]:
    """
    Nutzt doc.inline_shapes und den zugrundeliegenden Part-Blob fuer DOCX-Bilder.
    Speichert jedes Bild als JPEG + .meta.json, ruft Vision-API auf.
    """
    from docx import Document

    out_dir = _image_dir(path)
    out_dir.mkdir(parents=True, exist_ok=True)

    doc = Document(path)
    inline_shapes = doc.inline_shapes
    print(f"  DOCX: {len(inline_shapes)} Inline-Shapes", flush=True)

    metas: list[dict] = []

    for img_idx, ishape in enumerate(inline_shapes, 1):
        img_filename = f"img{img_idx}.jpg"
        img_path = out_dir / img_filename
        meta_path = out_dir / f"{img_filename}.meta.json"

        # Skip wenn bereits verarbeitet
        if _already_processed(meta_path):
            print(f"    Bild {img_idx}: bereits verarbeitet — ueberspringe", flush=True)
            existing = _read_meta(meta_path)
            if existing:
                metas.append(existing)
            continue

        # Blob aus dem Inline-Shape-Part lesen
        try:
            blob = ishape._inline.graphic.graphicData.pic.blipFill.blip.embed
            part = ishape.part.related_parts[blob]
            img_blob = part.blob
            from PIL import Image
            pil_img = Image.open(io.BytesIO(img_blob))
            jpeg_bytes = _to_jpeg_bytes(pil_img)
            if not jpeg_bytes:
                print(f"    Bild {img_idx}: Konvertierung fehlgeschlagen — ueberspringe", file=sys.stderr)
                continue
            width_px, height_px = _pil_size(Image.open(io.BytesIO(jpeg_bytes)))
        except Exception as e:
            print(f"    Bild {img_idx}: nicht lesbar ({e.__class__.__name__}) — ueberspringe", file=sys.stderr)
            continue

        # JPEG speichern
        img_path.write_bytes(jpeg_bytes)

        # Vision-API aufrufen (kein Seitenkontext in DOCX verfuegbar)
        print(f"    Bild {img_idx}: Vision-API...", flush=True)
        beschreibung = describe_image(jpeg_bytes, context="")

        # Meta speichern
        meta = {
            "source_file": _rel_to_wiki(path),
            "source_type": "docx",
            "bild_index": img_idx,
            "image_path": _rel_to_wiki(img_path),
            "beschreibung": beschreibung,
            "width_px": width_px,
            "height_px": height_px,
            "extracted_at": datetime.now(timezone.utc).isoformat(),
        }
        _write_meta(meta_path, meta)
        metas.append(meta)
        print(f"    Bild {img_idx}: gespeichert ({width_px}x{height_px})", flush=True)

    return metas


# ---------------------------------------------------------------------------
# Index aufbauen
# ---------------------------------------------------------------------------

def build_index() -> int:
    """
    Liest alle .meta.json aus CACHE_DIR und schreibt assets/image-index.json.
    Anschliessend wird wiki/picture_index.md als navigierbare Wiki-Seite geschrieben.
    Gibt die Anzahl der indizierten Bilder zurueck.
    """
    ASSETS_DIR.mkdir(parents=True, exist_ok=True)

    all_metas: list[dict] = []
    for meta_path in sorted(CACHE_DIR.rglob("*.meta.json")):
        meta = _read_meta(meta_path)
        if meta:
            all_metas.append(meta)

    INDEX_FILE.write_text(
        json.dumps(all_metas, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    print(f"Index geschrieben: {_rel_to_wiki(INDEX_FILE)} ({len(all_metas)} Bilder)", flush=True)

    # Wiki-Seite picture_index.md schreiben
    write_picture_index_md(all_metas)

    return len(all_metas)


# ---------------------------------------------------------------------------
# picture_index.md: navigierbare Wiki-Seite fuer den Bild-Index
# ---------------------------------------------------------------------------

def _lade_wiki_konzepte() -> list[str]:
    """
    Liest wiki/index.md und extrahiert alle [[wikilink]]-Bezeichner.
    Wird fuer das automatische Keyword-Matching verwendet.
    """
    index_path = WIKI_ROOT / "wiki" / "index.md"
    if not index_path.exists():
        return []
    import re
    inhalt = index_path.read_text(encoding="utf-8")
    # Alle [[bezeichner]] aus dem Index extrahieren
    return re.findall(r"\[\[([^\]]+)\]\]", inhalt)


def _konzepte_fuer_beschreibung(beschreibung: str, alle_konzepte: list[str]) -> list[str]:
    """
    Ermittelt passende Konzept-Links fuer eine Bildbeschreibung per Keyword-Matching.
    Sucht nach Woertern aus der Beschreibung (mind. 4 Zeichen) im Konzept-Bezeichner.
    Gibt maximal 3 passende [[konzept]]-Links zurueck.
    """
    if not beschreibung or not alle_konzepte:
        return []

    # Suchwoerter aus Beschreibung: Kleinbuchstaben, mind. 4 Zeichen
    woerter = set(
        w.lower().strip(".,;:!?()\"'")
        for w in beschreibung.split()
        if len(w) >= 4
    )

    treffer: list[str] = []
    for konzept in alle_konzepte:
        # Konzept-Bezeichner in Einzelteile zerlegen (Bindestriche als Trennzeichen)
        konzept_teile = set(konzept.lower().replace("-", " ").split())
        # Treffer wenn mindestens ein Suchwort im Konzept vorkommt
        if woerter & konzept_teile:
            treffer.append(f"[[{konzept}]]")
        if len(treffer) >= 3:
            break

    return treffer


def write_picture_index_md(all_metas: list[dict]) -> None:
    """
    Schreibt wiki/picture_index.md — eine navigierbare Wiki-Seite mit allen
    indizierten Bildern gruppiert nach Quelldatei.
    Konzept-Links werden per Keyword-Matching gegen wiki/index.md ermittelt.
    """
    wiki_dir = WIKI_ROOT / "wiki"
    wiki_dir.mkdir(parents=True, exist_ok=True)
    output_path = wiki_dir / "picture_index.md"

    heute = datetime.now(timezone.utc).strftime("%Y-%m-%d")

    # Konzepte aus Wiki-Index fuer Keyword-Matching laden
    alle_konzepte = _lade_wiki_konzepte()

    # Bilder nach Quelldatei gruppieren
    # Schluessel: Dateiname der Quelle (ohne Pfad-Praefix)
    from collections import defaultdict
    gruppen: dict[str, list[dict]] = defaultdict(list)
    for meta in all_metas:
        source = meta.get("source_file", "Unbekannt")
        # Nur den Dateinamen als Gruppenbezeichnung verwenden
        source_name = Path(source).name
        gruppen[source_name].append(meta)

    # Gesamtstatistik
    gesamt_bilder = len(all_metas)
    gesamt_quellen = len(gruppen)

    # Markdown aufbauen
    zeilen: list[str] = [
        "---",
        "title: Bild-Index",
        "type: index",
        f"last_updated: {heute}",
        "---",
        "",
        "# Bild-Index",
        "",
        "> Alle extrahierten Bilder aus dem Quell-Fundus — mit Beschreibung und Quellenangabe.",
        "> Navigiere von Bildern zu Konzepten und umgekehrt.",
        "",
        "## Überblick",
        "",
        f"- **Gesamt:** {gesamt_bilder} Bilder aus {gesamt_quellen} Quellen",
        f"- **Letzte Aktualisierung:** {heute}",
        "",
        "---",
        "",
        "## Bilder nach Quelle",
        "",
    ]

    # Jede Quelldatei als eigenen Abschnitt ausgeben
    for source_name in sorted(gruppen.keys()):
        eintraege = gruppen[source_name]
        zeilen.append(f"### {source_name}")
        zeilen.append("")
        zeilen.append("| Folie/Seite | Beschreibung | Konzepte |")
        zeilen.append("|-------------|-------------|---------|")

        for meta in eintraege:
            beschreibung = (meta.get("beschreibung") or "").replace("|", "\\|").replace("\n", " ")

            # Folien- oder Seiten-Nummer ermitteln
            if meta.get("source_type") == "pptx":
                position = str(meta.get("folie", "?"))
            elif meta.get("source_type") == "pdf":
                position = str(meta.get("seite", "?"))
            else:
                position = str(meta.get("bild_index", "?"))

            # Konzept-Links per Keyword-Matching ermitteln
            konzept_links = _konzepte_fuer_beschreibung(beschreibung, alle_konzepte)
            konzepte_str = ", ".join(konzept_links) if konzept_links else ""

            zeilen.append(f"| {position} | {beschreibung} | {konzepte_str} |")

        zeilen.append("")

    output_path.write_text("\n".join(zeilen), encoding="utf-8")
    print(f"picture_index.md geschrieben: wiki/picture_index.md ({gesamt_bilder} Bilder, {gesamt_quellen} Quellen)", flush=True)


# ---------------------------------------------------------------------------
# Verarbeitungs-Dispatcher
# ---------------------------------------------------------------------------

def process_file(path: Path) -> list[dict]:
    """
    Dispatcht die Extraktion an den richtigen Handler (PPTX/PDF/DOCX).
    Gibt Liste der extrahierten Meta-Dicts zurueck.
    Fehler werden geloggt aber nicht weitergeworfen.
    """
    suffix = path.suffix.lower()
    print(f"\nVerarbeite: {path.relative_to(WIKI_ROOT)}", flush=True)

    try:
        if suffix == ".pptx":
            return extract_images_pptx(path)
        elif suffix == ".pdf":
            return extract_images_pdf(path)
        elif suffix == ".docx":
            return extract_images_docx(path)
        else:
            print(f"  Kein Extraktor fuer '{suffix}' — ueberspringe", flush=True)
            return []
    except Exception as e:
        print(f"  FEHLER bei {path.name}: {e}", file=sys.stderr)
        return []


def process_all() -> None:
    """Verarbeitet alle PPTX/PDF/DOCX-Dateien unter RAW_DIR."""
    extensions = {".pptx", ".pdf", ".docx"}
    # Temporaere Office-Dateien (~$...) ueberspringen
    files = [
        p for p in sorted(RAW_DIR.rglob("*"))
        if p.suffix.lower() in extensions
        and not p.name.startswith("~$")
        and ".cache" not in p.parts
    ]

    print(f"Gefundene Dateien: {len(files)}", flush=True)
    total_images = 0

    for file_path in files:
        metas = process_file(file_path)
        total_images += len(metas)

    print(f"\nExtraktion abgeschlossen: {total_images} Bilder verarbeitet", flush=True)
    build_index()


# ---------------------------------------------------------------------------
# CLI-Einstiegspunkt
# ---------------------------------------------------------------------------

def main() -> None:
    args = sys.argv[1:]

    # --index-only: nur Index neu aufbauen
    if args == ["--index-only"]:
        build_index()
        return

    # Einzelne Datei angegeben
    if args and not args[0].startswith("--"):
        file_path = Path(args[0]).resolve()
        if not file_path.exists():
            print(f"Datei nicht gefunden: {file_path}", file=sys.stderr)
            sys.exit(1)
        process_file(file_path)
        build_index()
        return

    # Standardmodus: alle Dateien in raw/
    process_all()


if __name__ == "__main__":
    main()
